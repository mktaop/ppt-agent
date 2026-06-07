import streamlit as st
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pydantic import BaseModel, Field
import io, warnings, os
warnings.filterwarnings('ignore')


#--Pydantic Schemas
class ChartSeries(BaseModel):
    name: str = Field(description="Name of the data series (e.g., 'Revenue', 'Adoption Rate')")
    values: list[float] = Field(description="Numerical values for this series")

class ChartData(BaseModel):
    type: str = Field(description="Must be exactly one of: 'BAR', 'LINE', 'PIE'")
    categories: list[str] = Field(description="The X-axis labels or categories")
    series: list[ChartSeries] = Field(description="The data series to plot")
    summary: str = Field(description="A high-level summary of what this chart is conveying")

class Slide(BaseModel):
    is_title_slide: bool = Field(description="Set to True ONLY for the first slide.")
    header: str = Field(description="The main title of the slide.")
    content: list[str] = Field(description="Bullet points. If this is the title slide, put the subtitle here as a single string.")
    speaker_notes: str = Field(description="Speaker notes for this slide.")
    chart: ChartData | None = Field(default=None, description="Only include if a chart dramatically improves the explanation of the data.")

class PresentationDeck(BaseModel):
    slides: list[Slide]


def setup():
    st.set_page_config(page_title="AI Presentation Architect", 
                       page_icon="📊", 
                       layout="wide",
                       initial_sidebar_state="expanded",)
    
    hide_menu_style = """
            <style>
            #MainMenu {visibility: hidden;}
            </style>
            """
    st.markdown(hide_menu_style, unsafe_allow_html=True,)
    
    st.title("🚀 AI Presentation Agent")
    st.markdown("Choose to research a topic from the web, or upload your own documents to synthesize a custom deck.")

    
def get_template():
    st.sidebar.header("🎨 Theme & Design", divider='rainbow')
    st.sidebar.markdown("Upload a custom PowerPoint template (`.pptx`).")
    uploaded_template = st.sidebar.file_uploader("Upload Template (Optional)", type=["pptx"])
    return uploaded_template


def get_system_prompt():
    SYSTEM_PROMPT = """
                    You are an elite Presentation Architect. Your objective is to generate a logically ordered, fact-based, high-impact slide deck outline based on the user's inputs.
                    
                    CRITICAL INSTRUCTIONS:
                    1. Source Material: If the user provides a raw topic, you must use the Google Search tool to find the most up-to-date facts. If the user provides uploaded documents, you must extract and synthesize the most salient points directly from those documents.
                    2. Structure: Follow a logical narrative arc: Title Slide -> Hook -> Problem -> Solution -> Deep-Dive (Data/Facts) -> Conclusion.
                    3. Title Slide: Your first slide MUST be a title slide (`is_title_slide`: true).
                    4. Charts: When discussing trends, percentages, or comparisons, you MUST include a `chart` object. Use real data obtained from your search or the provided documents. Do not over-saturate; only use charts when visually appropriate.
                    """
    return SYSTEM_PROMPT


def create_pptx(deck: PresentationDeck, template_stream=None):
    if template_stream:
        prs = Presentation(template_stream)
    else:
        prs = Presentation()

    for slide_data in deck.slides:
        if slide_data.is_title_slide:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.header
            if slide_data.content:
                slide.placeholders[1].text = slide_data.content[0]
        else:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.header
            
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = "" 

            if slide_data.chart:
                TOP_MARGIN = Inches(1.5) 
                CHART_HEIGHT = Inches(4.5)
    
                body_shape.left = Inches(0.5)
                body_shape.top = TOP_MARGIN
                body_shape.width = Inches(4)
                body_shape.height = CHART_HEIGHT
                
                tf = body_shape.text_frame
                tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = f"📊 Key Insight: {slide_data.chart.summary}"
                p.font.bold = True
                p.font.size = Pt(16)
                
                for bullet in slide_data.content:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(14)
                
                chart_data_obj = CategoryChartData()
                chart_data_obj.categories = slide_data.chart.categories
                for series in slide_data.chart.series:
                    chart_data_obj.add_series(series.name, series.values)
                
                chart_type_map = {
                    'BAR': XL_CHART_TYPE.COLUMN_CLUSTERED,
                    'LINE': XL_CHART_TYPE.LINE,
                    'PIE': XL_CHART_TYPE.PIE
                }
                ctype = chart_type_map.get(slide_data.chart.type.upper(), XL_CHART_TYPE.COLUMN_CLUSTERED)
                
                x, y, cx, cy = Inches(4.8), TOP_MARGIN, Inches(4.5), CHART_HEIGHT
                slide.shapes.add_chart(ctype, x, y, cx, cy, chart_data_obj)
            else:
                for i, bullet in enumerate(slide_data.content):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = bullet
                    p.level = 0
        
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data.speaker_notes

    if template_stream and len(prs.slides) > len(deck.slides):
        try:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
        except Exception as e:
            st.warning(f"Could not remove template slide: {e}")

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io


def main():
    SYSTEM_PROMPT = get_system_prompt()
    uploaded_template = get_template()
    
    tab1, tab2 = st.tabs(["🌐 Web Research", "📄 Document Upload"])
    
    contents_to_send = []
    tools_to_use = None
    ready_to_generate = False
    filename_prefix = "presentation"
    
    with tab1:
        st.markdown("### Generate from Web Research")
        topic = st.text_input("What is the presentation topic?", placeholder="e.g., Global EV Market Share 2024 vs 2026", key="web_topic")
        
        if st.button("Generate from Web", type="primary", key="btn_web"):
            if not topic:
                st.warning("Please enter a topic first.")
            else:
                contents_to_send = [f"Create a presentation about: {topic}"]
                tools_to_use = [types.Tool(google_search=types.GoogleSearch())]
                filename_prefix = topic.replace(' ', '_')
                ready_to_generate = True
    
    with tab2:
        st.markdown("### Generate from PDFs")
        uploaded_pdfs = st.file_uploader("Upload one or more PDFs", type=["pdf"], accept_multiple_files=True)
        pdf_focus = st.text_input("Optional: What specific aspects of these documents should I focus on?", placeholder="e.g., Focus only on the Q3 financial results.")
        
        if st.button("Generate from Documents", type="primary", key="btn_doc"):
            if not uploaded_pdfs:
                st.warning("Please upload at least one PDF.")
            else:
                # Append Native PDF Parts to the Gemini prompt
                for pdf in uploaded_pdfs:
                    contents_to_send.append(
                        types.Part.from_bytes(
                            data=pdf.getvalue(),
                            mime_type='application/pdf'
                        )
                    )
                
                # Add the text directive
                if pdf_focus:
                    contents_to_send.append(f"Extract the most salient information from these documents to build a presentation, focusing specifically on: {pdf_focus}")
                else:
                    contents_to_send.append("Extract the most important, salient information from these documents to build a comprehensive presentation.")
                
                tools_to_use = None # No search needed, rely on docs
                filename_prefix = "Document_Synthesis"
                ready_to_generate = True
    
    if ready_to_generate:
        with st.status("Agent at work...", expanded=True) as status:
            st.write("🧠 Analyzing inputs and designing slide structure...")
            
            try:
                # Build Config dynamically based on mode
                config_params = {
                    "system_instruction": SYSTEM_PROMPT,
                    "response_mime_type": "application/json",
                    "response_schema": PresentationDeck, 
                    "temperature": 0.2,
                    "thinking_config": types.ThinkingConfig(include_thoughts=True),
                }
                if tools_to_use:
                    config_params["tools"] = tools_to_use
    
                config = types.GenerateContentConfig(**config_params)
    
                # API Call
                response = client.models.generate_content(
                    model=MODEL_ID,
                    contents=contents_to_send,
                    config=config
                )
                
                deck: PresentationDeck = response.parsed
                
                st.write("🎨 Assembling PowerPoint file...")
                template_stream = io.BytesIO(uploaded_template.getvalue()) if uploaded_template else None
                pptx_file = create_pptx(deck, template_stream=template_stream)
                
                status.update(label="Presentation Ready!", state="complete", expanded=False)
                st.success("Your presentation has been synthesized and built.")
                
                st.download_button(
                    label="📥 Download PowerPoint",
                    data=pptx_file,
                    file_name=f"{filename_prefix}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
    
                with st.expander("View Outline Preview", expanded=True):
                    for idx, slide in enumerate(deck.slides):
                        st.markdown(f"### Slide {idx + 1}: {slide.header}")
                        if slide.is_title_slide:
                            st.caption("(Title Slide)")
                        for bullet in slide.content:
                            st.markdown(f"- {bullet}")
                        if slide.chart:
                            st.info(f"📈 **Includes {slide.chart.type} Chart:** {slide.chart.summary}")
                        st.markdown("---")
    
            except Exception as e:
                status.update(label="Error occurred", state="error", expanded=False)
                st.error(f"An error occurred during generation: {e}")


if __name__ == '__main__':
    try:
        GEMINI_API_KEY = os.environ.get("GOOGLE_API_KEY")
        client = genai.Client(api_key=GEMINI_API_KEY)
        MODEL_ID = "gemini-3-flash-preview"
    except Exception:
        st.error("🚨 GOOGLE_API_KEY not found in environment variables.")
        st.stop()
    
    setup()
    main()
